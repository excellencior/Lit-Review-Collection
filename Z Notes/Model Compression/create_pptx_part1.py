"""Gemma 4 Efficiency Presentation - Helper Functions"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Academic color palette
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x6B)
ACCENT = RGBColor(0x3A, 0x7C, 0xA5)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x6B, 0x7B, 0x8D)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED_ACC = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
TABLE_HEAD = RGBColor(0x2C, 0x3E, 0x6B)
TABLE_ALT = RGBColor(0xEB, 0xEF, 0xF5)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)  # 16:9

def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, w, h, text, font_size=18, color=DARK_TEXT,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def set_para(tf, text, size=18, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT,
             space_before=0, space_after=0, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p

def make_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, NAVY)
    # Top accent line
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)
    add_text_box(slide, Inches(1.2), Inches(2.2), Inches(10.9), Inches(1.8),
                 title, 40, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.2), Inches(4.0), Inches(10.9), Inches(1.0),
                 subtitle, 20, RGBColor(0xA0,0xB4,0xCC), False, PP_ALIGN.CENTER)
    return slide

def make_section_slide(prs, number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BLUE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)
    add_text_box(slide, Inches(1.2), Inches(2.5), Inches(10.9), Inches(0.8),
                 f"SECTION {number}", 16, RGBColor(0x7F,0x9C,0xBF), True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.2), Inches(3.2), Inches(10.9), Inches(1.2),
                 title, 36, WHITE, True, PP_ALIGN.CENTER)
    return slide

def slide_header(slide, title):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)
    add_rect(slide, Inches(0), Inches(0.06), SLIDE_W, Inches(1.0), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.18), Inches(11.7), Inches(0.7),
                 title, 26, WHITE, True, PP_ALIGN.LEFT)

def content_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    slide_header(slide, title)
    return slide

def add_table(slide, rows, cols, left, top, w, h, data, col_widths=None):
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, w, h)
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = DARK_TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEAD
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return tbl_shape

def bullet_block(slide, left, top, w, h, items, title=None, title_size=18, item_size=15):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(item_size)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Calibri"
        p.space_before = Pt(3)
        p.space_after = Pt(3)
    return txBox
